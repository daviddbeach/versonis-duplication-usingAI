"""
pii_engine.py — shared PII detection, scoring, masking, and text-extraction
helpers used by both functions of the Data Privacy skill.

Detection stack:
  * Microsoft Presidio (AnalyzerEngine) for the heavy lifting.
  * spaCy NER (en_core_web_lg preferred, en_core_web_sm fallback) for
    unstructured PII: PERSON names and LOCATION/address fragments.
  * Regex pattern recognizers for structured PII. Presidio ships strong
    built-ins (email, phone via libphonenumber, credit card w/ Luhn,
    US SSN, passport, IBAN, IP, etc.); we add a Canadian SIN recognizer
    (Luhn-validated) since the user base spans Canada/US.
  * Tesseract OCR (via pytesseract) to read PII out of images.

Nothing in here moves, deletes, or overwrites files. File mutation lives
in the two function scripts, behind an explicit confirmation gate.
"""

from __future__ import annotations

import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

# ---------------------------------------------------------------------------
# Sensitivity model
# ---------------------------------------------------------------------------
# HIGH-tier entities are sensitive enough that a single confident hit marks a
# file as a "personal information file". LOW-tier entities (a lone name, one
# email) are too common in ordinary business documents to quarantine on their
# own; they must co-occur or reach a density threshold. See score_file().

HIGH_TIER = {
    "US_SSN",
    "CA_SIN",
    "CREDIT_CARD",
    "US_PASSPORT",
    "US_DRIVER_LICENSE",
    "US_BANK_NUMBER",
    "IBAN_CODE",
    "US_ITIN",
    "MEDICAL_LICENSE",
    "CRYPTO",
}

LOW_TIER = {
    "PERSON",
    "EMAIL_ADDRESS",
    "PHONE_NUMBER",
    "LOCATION",
    "IP_ADDRESS",
    "DATE_TIME",
    "NRP",
}

# Entities we ask Presidio for. Order doesn't matter.
DEFAULT_ENTITIES = sorted(HIGH_TIER | LOW_TIER)

# Sensitivity presets control how eagerly a file is flagged in Function 1.
SENSITIVITY_PRESETS = {
    # only HIGH-tier identifiers trip the wire — fewest false positives
    "low": {"high_only": True, "min_distinct_low": 99, "min_low_count": 999},
    # default: HIGH hit, OR 2+ distinct LOW types, OR 4+ LOW instances
    "medium": {"high_only": False, "min_distinct_low": 2, "min_low_count": 4},
    # aggressive: any single PII hit of any tier flags the file
    "high": {"high_only": False, "min_distinct_low": 1, "min_low_count": 1},
}

# File extensions grouped by how we pull text out of them.
PLAINTEXT_EXTS = {".txt", ".md", ".markdown", ".csv", ".tsv", ".json",
                  ".log", ".yaml", ".yml", ".xml", ".html", ".htm", ".rtf"}
DOCX_EXTS = {".docx"}
XLSX_EXTS = {".xlsx", ".xlsm"}
PPTX_EXTS = {".pptx"}
PDF_EXTS = {".pdf"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif", ".webp"}

SCANNABLE_EXTS = (PLAINTEXT_EXTS | DOCX_EXTS | XLSX_EXTS | PPTX_EXTS
                  | PDF_EXTS | IMAGE_EXTS)


class ModelNotInstalled(RuntimeError):
    """Raised when no spaCy NER model is available."""


# ---------------------------------------------------------------------------
# Analyzer construction
# ---------------------------------------------------------------------------
def _canadian_sin_recognizer():
    """Canadian Social Insurance Number, Luhn-validated to cut false hits."""
    from presidio_analyzer import Pattern, PatternRecognizer

    pattern = Pattern(
        name="ca_sin",
        regex=r"\b\d{3}[-\s]?\d{3}[-\s]?\d{3}\b",
        score=0.3,  # boosted to 0.7 below when Luhn passes
    )

    class SINRecognizer(PatternRecognizer):
        def validate_result(self, pattern_text: str):
            digits = [int(c) for c in pattern_text if c.isdigit()]
            if len(digits) != 9:
                return None
            total = 0
            for i, d in enumerate(digits):
                if i % 2 == 1:  # double every second digit
                    d *= 2
                    if d > 9:
                        d -= 9
                total += d
            return total % 10 == 0  # True -> score bumped, False -> dropped

    return SINRecognizer(
        supported_entity="CA_SIN",
        patterns=[pattern],
        context=["sin", "social insurance", "numéro d'assurance"],
    )


def _blank_nlp_engine():
    """A Presidio NLP engine backed by a blank spaCy pipeline (NO names).

    Used as a degraded fallback when no NER model is installed. Structured
    detectors (email, SSN, card, SIN, phone, IP) still work; PERSON/LOCATION
    will be missed. Gated by the caller behind a clear warning.
    """
    import spacy
    from presidio_analyzer.nlp_engine import SpacyNlpEngine
    from presidio_analyzer.nlp_engine.ner_model_configuration import (
        NerModelConfiguration,
    )

    nlp = spacy.blank("en")
    if "sentencizer" not in nlp.pipe_names:
        nlp.add_pipe("sentencizer")

    class _Loaded(SpacyNlpEngine):
        def __init__(self, loaded):
            self.nlp = {"en": loaded}
            self.ner_model_configuration = NerModelConfiguration()

    return _Loaded(nlp)


def get_nlp_engine(preferred_models=("en_core_web_lg", "en_core_web_sm")):
    """Return a Presidio NlpEngine backed by the best available spaCy model.

    If no model is installed: raise ModelNotInstalled, UNLESS the env var
    DP_ALLOW_BLANK_NLP is set, in which case fall back to a blank pipeline
    (structured PII only) with a warning so the run can still proceed.
    """
    import spacy
    from presidio_analyzer.nlp_engine import NlpEngineProvider

    for model_name in preferred_models:
        if spacy.util.is_package(model_name):
            provider = NlpEngineProvider(nlp_configuration={
                "nlp_engine_name": "spacy",
                "models": [{"lang_code": "en", "model_name": model_name}],
            })
            return provider.create_engine(), model_name

    if os.environ.get("DP_ALLOW_BLANK_NLP"):
        print("WARNING: no spaCy NER model installed — running in degraded "
              "mode. Names/addresses will NOT be detected. Install "
              "en_core_web_lg via the setup step for full coverage.",
              file=__import__("sys").stderr)
        return _blank_nlp_engine(), "blank (degraded: no name detection)"

    raise ModelNotInstalled(
        "No spaCy English model found. Run the skill's setup step "
        "(scripts/setup_env.py) to install en_core_web_lg."
    )


def build_analyzer(nlp_engine=None):
    """Build a Presidio AnalyzerEngine with our extra recognizers.

    Pass a prebuilt nlp_engine (used by tests with a blank pipeline);
    otherwise the best installed spaCy model is loaded.
    """
    from presidio_analyzer import AnalyzerEngine

    model_name = None
    if nlp_engine is None:
        nlp_engine, model_name = get_nlp_engine()

    analyzer = AnalyzerEngine(nlp_engine=nlp_engine, supported_languages=["en"])
    analyzer.registry.add_recognizer(_canadian_sin_recognizer())
    analyzer._dp_model_name = model_name  # informational, for logging
    return analyzer


# ---------------------------------------------------------------------------
# Analysis + scoring
# ---------------------------------------------------------------------------
@dataclass
class FileVerdict:
    path: str
    flagged: bool
    reason: str
    entity_counts: dict = field(default_factory=dict)
    high_hits: list = field(default_factory=list)
    error: Optional[str] = None


def analyze_text(analyzer, text: str, min_score: float = 0.4, entities=None):
    """Run Presidio; keep results at or above min_score."""
    if not text or not text.strip():
        return []
    results = analyzer.analyze(text=text, language="en",
                               entities=entities or DEFAULT_ENTITIES)
    return [r for r in results if r.score >= min_score]


def score_file(path: str, results, *, sensitivity: str = "medium") -> FileVerdict:
    """Decide whether a file is a personal-information file.

    Uses the sensitivity preset's thresholds over the detected entities.
    """
    cfg = SENSITIVITY_PRESETS[sensitivity]
    counts: dict[str, int] = {}
    for r in results:
        counts[r.entity_type] = counts.get(r.entity_type, 0) + 1

    high_hits = [e for e in counts if e in HIGH_TIER]
    low_types = [e for e in counts if e in LOW_TIER]
    low_total = sum(counts[e] for e in low_types)

    if high_hits:
        return FileVerdict(path, True,
                           f"high-sensitivity: {', '.join(sorted(high_hits))}",
                           counts, high_hits)

    if cfg["high_only"]:
        return FileVerdict(path, False, "no high-sensitivity PII", counts, [])

    if len(low_types) >= cfg["min_distinct_low"]:
        return FileVerdict(path, True,
                           f"{len(low_types)} distinct PII types: "
                           f"{', '.join(sorted(low_types))}", counts, [])

    if low_total >= cfg["min_low_count"]:
        return FileVerdict(path, True,
                           f"{low_total} PII instances "
                           f"({', '.join(sorted(low_types))})", counts, [])

    return FileVerdict(path, False, "below threshold", counts, [])


# ---------------------------------------------------------------------------
# Masking
# ---------------------------------------------------------------------------
def mask_text(text: str, results, *, style: str = "asterisks",
              mask: str = "*****"):
    """Return text with every detected PII span replaced.

    style="asterisks" -> every span becomes `mask` (default "*****").
    style="tokens"    -> span becomes [REDACTED_<ENTITY>], preserving the
                         kind of value so downstream models keep structure.
    """
    from presidio_anonymizer import AnonymizerEngine
    from presidio_anonymizer.entities import OperatorConfig

    engine = AnonymizerEngine()

    if style == "tokens":
        entity_types = {r.entity_type for r in results}
        operators = {
            et: OperatorConfig("replace",
                               {"new_value": f"[REDACTED_{et}]"})
            for et in entity_types
        }
        operators["DEFAULT"] = OperatorConfig("replace",
                                              {"new_value": "[REDACTED]"})
    else:
        operators = {"DEFAULT": OperatorConfig("replace",
                                               {"new_value": mask})}

    out = engine.anonymize(text=text, analyzer_results=results,
                           operators=operators)
    return out.text


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------
def extract_text(path: str) -> str:
    """Best-effort plain-text extraction for any supported file type.

    Raises on unsupported types so the caller can skip + log them.
    """
    p = Path(path)
    ext = p.suffix.lower()

    if ext in PLAINTEXT_EXTS:
        return p.read_text(encoding="utf-8", errors="ignore")

    if ext in DOCX_EXTS:
        import docx
        d = docx.Document(str(p))
        parts = [para.text for para in d.paragraphs]
        for table in d.tables:
            for row in table.rows:
                parts.extend(cell.text for cell in row.cells)
        return "\n".join(parts)

    if ext in XLSX_EXTS:
        import openpyxl
        wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                parts.append(" ".join(str(c) for c in row if c is not None))
        wb.close()
        return "\n".join(parts)

    if ext in PPTX_EXTS:
        from pptx import Presentation
        prs = Presentation(str(p))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
        return "\n".join(parts)

    if ext in PDF_EXTS:
        import pdfplumber
        parts = []
        with pdfplumber.open(str(p)) as pdf:
            for page in pdf.pages:
                parts.append(page.extract_text() or "")
        text = "\n".join(parts)
        # Scanned PDFs yield ~no text; flag that to the caller via marker.
        if len(text.strip()) < 10:
            return "\n__DP_SCANNED_PDF__\n" + text
        return text

    if ext in IMAGE_EXTS:
        import pytesseract
        from PIL import Image
        return pytesseract.image_to_string(Image.open(str(p)))

    raise ValueError(f"unsupported file type: {ext}")
